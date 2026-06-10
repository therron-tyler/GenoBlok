from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pathlib import Path

BASE = Path("/path/to/data/IMPACT_ClusteringRevisited")
OUT = BASE / "IMPACT_ClusteringRevisited.pptx"

SUBFOLDERS = ["dims_1-10_plots", "dims_1-11_plots", "dims_1-12_plots", "dims_1-9_11-12_plots"]

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]  # completely blank

for folder_name in SUBFOLDERS:
    folder = BASE / folder_name
    # Extract dimension string from folder name e.g. "dims_1-10_plots" -> "dims 1-10"
    dim_label = folder_name.replace("_plots", "").replace("_", " ")

    images = sorted(folder.glob("*.png"))
    for img_path in images:
        # Clean filename stem into a plain title, e.g. "03a_umap_clusters" -> "umap clusters"
        stem = img_path.stem
        parts = stem.split("_", 1)
        plot_name = parts[1] if len(parts) == 2 and parts[0][0].isdigit() else stem
        plot_name = plot_name.replace("_", " ")

        slide = prs.slides.add_slide(blank_layout)

        # Title text box at the top
        title_box = slide.shapes.add_textbox(Inches(0.2), Inches(0.1), Inches(12.9), Inches(0.5))
        tf = title_box.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{plot_name}  |  {dim_label}  |  res 0.3"
        run.font.size = Pt(18)
        run.font.bold = False

        # Image fills remaining slide area
        slide.shapes.add_picture(
            str(img_path),
            Inches(0.2), Inches(0.65),
            width=Inches(12.9), height=Inches(6.7)
        )

prs.save(OUT)
print(f"Saved: {OUT}  ({len(prs.slides)} slides)")
