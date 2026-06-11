#!/usr/bin/env python3
"""Build a simple PowerPoint from the NPSLE PCA/SCT investigation pipeline outputs.

Organizes plots in the order the R pipeline produces them:
  1. Pipeline goal / overview
  2. Pre-SCT PCA (LogNormalize): elbow, per-PC-pair drivers + biplots, DimHeatmap, violins
  3. Post-SCT PCA (SCTransform v2): same layout
  4. Cluster sweep at res=0.4 across dims 6..10: UMAP, sample composition, marker heatmap
"""

from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN

PLOT_DIR = Path("/path/to/data/20260526_NPSLE_PCA_SCT_Tcells_plots")
OUT_PATH = Path("/path/to/data/20260526_NPSLE_PCA_SCT_Tcells_deck.pptx")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]


def add_text_slide(title, body_lines):
    slide = prs.slides.add_slide(BLANK)
    tb = slide.shapes.add_textbox(Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.9))
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    body = slide.shapes.add_textbox(Inches(0.5), Inches(1.4), Inches(12.3), Inches(5.6))
    btf = body.text_frame
    btf.word_wrap = True
    for i, line in enumerate(body_lines):
        para = btf.paragraphs[0] if i == 0 else btf.add_paragraph()
        para.text = line
        para.font.size = Pt(16)
        para.space_after = Pt(6)


def add_image_slide(title, caption, image_name):
    """One image on a slide with a title bar at the top and a short caption below."""
    slide = prs.slides.add_slide(BLANK)

    # Title
    tb = slide.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.7), Inches(0.55))
    tp = tb.text_frame.paragraphs[0]
    tp.text = title
    tp.font.size = Pt(20)
    tp.font.bold = True

    # Caption
    cap = slide.shapes.add_textbox(Inches(0.3), Inches(0.72), Inches(12.7), Inches(0.5))
    cp = cap.text_frame.paragraphs[0]
    cap.text_frame.word_wrap = True
    cp.text = caption
    cp.font.size = Pt(12)
    cp.font.italic = True

    # Image — fit within remaining area (about 12.7w x 6.0h)
    img_path = PLOT_DIR / image_name
    if not img_path.exists():
        msg = slide.shapes.add_textbox(Inches(0.5), Inches(3.0), Inches(12.3), Inches(1.0))
        msg.text_frame.paragraphs[0].text = f"[MISSING] {image_name}"
        return

    # Place with a sensible max size; pptx will preserve aspect by giving only width.
    max_w = Inches(12.5)
    max_h = Inches(6.0)
    pic = slide.shapes.add_picture(str(img_path), Inches(0.4), Inches(1.3),
                                    width=max_w)
    # If too tall, rescale by height
    if pic.height > max_h:
        ratio = pic.width / pic.height
        pic.height = max_h
        pic.width = int(max_h * ratio)
    # Recentre horizontally
    pic.left = int((prs.slide_width - pic.width) / 2)


# ---------------------------------------------------------------------------
# 1. Title + goal slides
# ---------------------------------------------------------------------------
add_text_slide(
    "NPSLE T-cell PCA / SCTransform Investigation",
    [
        "Pipeline: pca_driver_investigation.R",
        "Plots: 20260526_NPSLE_PCA_SCT_Tcells_plots",
        "",
        "Purpose: identify which genes drive the principal components in the NPSLE T-cell "
        "subcluster before vs. after SCTransform v2, then sweep clustering across PCA "
        "dimensionalities to see which (dims, resolution) yields biologically clean clusters.",
        "",
        "Per-sample resolution: every plot is coloured by HTO_maxID so it is obvious when a "
        "PC or cluster is being driven by a single donor rather than shared biology.",
    ],
)

add_text_slide(
    "Why this analysis exists",
    [
        "1. PCs in scRNA-seq are easily hijacked by technical noise (mito/ribo, TCR/Ig V-genes, "
        "cell cycle, single-donor batch effects). The pipeline strips those gene families from "
        "the HVG set and regresses percent.mt + S.Score + G2M.Score so PCs reflect biology.",
        "",
        "2. SCTransform v2 (Pearson residuals) is expected to normalize sequencing-depth bias "
        "better than LogNormalize. We compare LogNormalize-PCA (\"pre-SCT\") vs SCT-PCA "
        "side-by-side to confirm SCT removed depth/technical structure.",
        "",
        "3. Once we trust the embedding, the cluster sweep (dims = 6,7,8,9,10 at resolution 0.4) "
        "explores how cluster count and biology change with PCA dimensionality, with de novo "
        "marker heatmaps to judge each candidate.",
    ],
)

add_text_slide(
    "Plot vocabulary used throughout",
    [
        "- Elbow plot: variance per PC; the \"elbow\" suggests how many PCs carry signal.",
        "- Drivers (scatter + barplots): top scatter shows cells in the PC plane coloured by "
        "sample; bottom diverging bars show the top +/- loading genes for each PC (red = "
        "positive loading, blue = negative).",
        "- Biplot: cells coloured by sample with arrows showing the strongest gene-loading "
        "vectors. Direction of an arrow tells you which PC region that gene pushes cells into.",
        "- DimHeatmap: cells (columns, ranked by PC score) x top loading genes (rows). A clean "
        "diagonal block pattern means the PC is driven by a coherent gene set.",
        "- Per-sample PC-score violins: distribution of each PC's score per HTO sample. Lets "
        "you spot a PC that is really a single-donor axis.",
        "- Cluster UMAP / sample composition / de novo marker heatmap: standard outputs of "
        "FindClusters -> FindAllMarkers used to judge each (dim, res) combination.",
    ],
)

# ---------------------------------------------------------------------------
# 2. Pre-SCT PCA stage
# ---------------------------------------------------------------------------
add_text_slide(
    "Stage 1 - Pre-SCT PCA (LogNormalize baseline)",
    [
        "Reduction: pca_lognorm  (keys: PClog_1, PClog_2, ...)",
        "",
        "Pipeline: NormalizeData -> CellCycleScoring -> FindVariableFeatures(3000) "
        "-> strip MT/ribo/TCR/Ig V-genes -> ScaleData (regress percent.mt + S + G2M) -> RunPCA.",
        "",
        "Purpose of this section: establish the baseline PCA before SCT correction. Any "
        "single-sample axes or depth-driven loadings we see here are exactly what SCT should "
        "minimize in Stage 2.",
    ],
)

add_image_slide(
    "Pre-SCT: Elbow plot",
    "Variance explained per PC after LogNormalize. The elbow tells you the rough number of "
    "informative PCs and lets you compare against the post-SCT elbow.",
    "00_elbow_pca_lognorm.png",
)

for d1, d2 in [(1, 2), (3, 4), (5, 6), (7, 8), (9, 10)]:
    add_image_slide(
        f"Pre-SCT: PClog_{d1} vs PClog_{d2} - scores + driver genes",
        "Top: cells in the PC plane coloured by HTO_maxID (single-sample dominance is a red flag). "
        "Bottom: diverging top +/- loading genes per PC - these genes literally define the axis.",
        f"presct_PClog_{d1}_PClog_{d2}_drivers.png",
    )
    add_image_slide(
        f"Pre-SCT: PClog_{d1} vs PClog_{d2} - biplot",
        "Cells coloured by sample with the strongest gene-loading vectors overlaid as arrows. "
        "Arrow direction shows which corner of the PC plane that gene pulls cells toward.",
        f"presct_PClog_{d1}_PClog_{d2}_biplot.png",
    )

add_image_slide(
    "Pre-SCT: DimHeatmap PC1-10",
    "Cells (columns, ranked by each PC's score) x top loading genes (rows), one panel per PC. "
    "A clean two-block diagonal means the PC is a coherent gene program; a noisy block means "
    "the PC is weak or technical.",
    "presct_dimheatmap_PC1_10.png",
)

add_image_slide(
    "Pre-SCT: Per-sample PC-score violins",
    "Distribution of PC scores per HTO sample for PC1-PC10. Any PC where one or two samples "
    "sit completely apart from the rest is a sample-driven axis, not biology - this is the "
    "single most useful sanity check before SCT.",
    "presct_persample_pc_score_violins.png",
)

# ---------------------------------------------------------------------------
# 3. Post-SCT PCA stage
# ---------------------------------------------------------------------------
add_text_slide(
    "Stage 2 - Post-SCT PCA (SCTransform v2)",
    [
        "Reduction: pca_sct  (keys: PCsct_1, PCsct_2, ...)",
        "",
        "Pipeline: SCTransform v2 (glmGamPoi, regress percent.mt + S + G2M, 3000 HVGs) "
        "-> strip MT/ribo/TCR/Ig V-genes from HVG list -> RunPCA on the SCT assay.",
        "",
        "Read this section directly against Stage 1: SCT should (a) flatten the elbow into a "
        "smoother decay, (b) replace single-donor PCs with shared T-cell biology, and (c) put "
        "interpretable gene programs on the leading PCs.",
    ],
)

add_image_slide(
    "Post-SCT: Elbow plot",
    "Variance explained per PC on the SCT-transformed data. Compare elbow location and shape to "
    "the pre-SCT elbow to judge how much depth-driven variance SCT absorbed.",
    "01_elbow_pca_sct.png",
)

for d1, d2 in [(1, 2), (3, 4), (5, 6), (7, 8), (9, 10)]:
    add_image_slide(
        f"Post-SCT: PCsct_{d1} vs PCsct_{d2} - scores + driver genes",
        "Same layout as the pre-SCT drivers plot. Look for: (i) less single-sample colour "
        "stratification in the scatter, (ii) recognizable T-cell programs (cytotoxicity, "
        "interferon, naive/memory, etc.) in the top loading bars.",
        f"sct_PCsct_{d1}_PCsct_{d2}_drivers.png",
    )
    add_image_slide(
        f"Post-SCT: PCsct_{d1} vs PCsct_{d2} - biplot",
        "Biplot on the SCT embedding. Arrows now point along biology-driven axes if SCT worked; "
        "if they still point along sample-coloured spreads, that PC is still confounded.",
        f"sct_PCsct_{d1}_PCsct_{d2}_biplot.png",
    )

add_image_slide(
    "Post-SCT: DimHeatmap PC1-10",
    "DimHeatmap on the SCT assay. Cleaner diagonals than the pre-SCT version indicate SCT "
    "produced more coherent PCs - this is the visual confirmation that depth-correction worked.",
    "sct_dimheatmap_PC1_10.png",
)

add_image_slide(
    "Post-SCT: Per-sample PC-score violins",
    "Per-sample distributions on PCsct_1..10. The goal here is OVERLAPPING violins across "
    "samples - that means the PC captures shared biology rather than a donor effect. PCs where "
    "one sample is still offset are candidates to drop from downstream dims.",
    "sct_persample_pc_score_violins.png",
)

# ---------------------------------------------------------------------------
# 4. Cluster sweep
# ---------------------------------------------------------------------------
add_text_slide(
    "Stage 3 - Cluster sweep on pca_sct (resolution = 0.4)",
    [
        "FindNeighbors -> FindClusters -> RunUMAP on pca_sct for each dim in {6, 7, 8, 9, 10}.",
        "PrepSCTFindMarkers then FindAllMarkers (Wilcoxon, only.pos, log2FC > 0.25, padj < 0.05); "
        "top 5 markers per cluster shown on the de novo heatmap.",
        "",
        "Each dim gets three plots in this order:",
        "  - UMAP (clusters | sample) - shape of the embedding + sample mixing",
        "  - Sample composition - % of each cluster contributed by each HTO sample",
        "  - De novo marker heatmap - top 5 markers per cluster, scaled expression",
        "",
        "How to read this section: walk d6 -> d10 and ask which dims give clusters that are "
        "(i) sample-mixed, (ii) marker-distinct, (iii) not over-split.",
    ],
)

for d in [6, 7, 8, 9, 10]:
    add_image_slide(
        f"Cluster sweep d={d}, res=0.4 - UMAP",
        "Left: UMAP coloured by Leiden cluster ID. Right: same UMAP coloured by HTO sample. "
        "Cluster islands that match donor colour 1:1 are donor effects, not cell states.",
        f"cluster_d{d}_res0.4_umap.png",
    )
    add_image_slide(
        f"Cluster sweep d={d}, res=0.4 - sample composition",
        "Stacked bars: each cluster's makeup by HTO sample (% of cluster). A cluster dominated "
        "by one sample is suspect; clusters drawing roughly proportionally from all samples are "
        "the credible cell states.",
        f"cluster_d{d}_res0.4_sample_composition.png",
    )
    add_image_slide(
        f"Cluster sweep d={d}, res=0.4 - de novo marker heatmap",
        "Cells (columns, grouped by cluster) x top 5 positively-enriched markers per cluster "
        "(rows). Sharp on-diagonal blocks of high expression = clusters with distinct identities. "
        "Bleed-across or noisy blocks suggest the clustering is over-splitting or the dim choice "
        "is wrong.",
        f"cluster_d{d}_res0.4_denovo_heatmap.png",
    )

# ---------------------------------------------------------------------------
# Save
# ---------------------------------------------------------------------------
prs.save(OUT_PATH)
print(f"Wrote: {OUT_PATH}")
print(f"Slides: {len(prs.slides)}")
